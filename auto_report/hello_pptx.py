        # -*- coding: utf-8 -*-
"""
Created on Thu Apr  5 12:12:34 2018

@author: Administrator
"""

# 数据分析
import matplotlib.pyplot as plt
import pandas as pd

# PPT
from pptx.util import Inches 

# slide layout测试
def add_all_slide_layout(prs):
    slide_layout_list = ['Title (presentation title slide)',
                         'Title and Content',
                         'Section Header (sometimes called Segue)',
                         'Two Content (side by side bullet textboxes)',
                         'Comparison (same but additional title for each side by side content box)',
                         'Title Only',
                         'Blank',
                         'Content with Caption',
                         'Picture with Caption']
    
    for index, name in enumerate(slide_layout_list):
        new_slide_layout = prs.slide_layouts[index]
        new_slide = prs.slides.add_slide(new_slide_layout)
        
        #6对应的是Blank，没有标题框
        if index != 6:
            title = new_slide.shapes.title
            title.text = name
            
# slide shape测试
def add_text_by_shape(slide):
    # 遍历所有的shape
    for i, shape in enumerate(slide.shapes):
        # 打印shape的id和名称
        print('%d %s' % (shape.shape_id, shape.name))
        
        # 如果没有文本框就跳过
        if not shape.has_text_frame:
            continue
        
        # 设置shape的文本
        shape.text = 'add_text_by_shape ' + str(i)

def add_text_by_palceholder(slide):
    for i, shape in enumerate(slide.placeholders):
        print('%d %s' % (shape.placeholder_format.idx, shape.name))
        if not shape.has_text_frame:
            continue
        
        shape.text = 'add_text_by_palceholder ' + str(i)  

# 箱线图    
def page1_boxplot(prs):
    # slide_layouts[1]为带标题和正文框的ppt
    title_slide_layout = prs.slide_layouts[1]
    
    # 新增ppt（当前仅支持新增操作）
    slide = prs.slides.add_slide(title_slide_layout)
    
    # 取本页ppt的title，向title文本框写如文字 
    title = slide.shapes.title
    title.text = '箱线图示例' 
    
    # 导入数据集
    df = pd.read_excel('../resources/data/data0.xlsx', 'Sheet1')
    
    # 绘制箱线图 
    # 方法1：使用matplotlib，暂未找到方法，待完善
    # fig = plt.figure()
    # ax = fig.add_subplot(111)
    # ax.boxplot(df['Lead Time'])
    
    # 方法2：使用dataframe，第一个参数column是数据列名，第二个参数by是分组的列名
    df.boxplot(column='Lead Time', by='Project')
    
    plt.savefig('../resources/report/交付周期_箱线图.png') 
    plt.show()
    
    # 插入图片，在指定位置按预设值添加图片 
    img_path = '../resources/report/交付周期_箱线图.png' 
    left, top, width, height = Inches(1), Inches(1.8), Inches(4.5), Inches(4.5)  
    slide.shapes.add_picture(img_path, left, top, width, height)
    